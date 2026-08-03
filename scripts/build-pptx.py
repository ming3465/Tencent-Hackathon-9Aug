"""Legacy fallback exporter for the raster-source judge deck.

The 2026-08-03 campaign artifact was edited and verified with artifact-tool.
This script remains a project-local fallback for teammates who already use
python-pptx; it preserves the same eight 16:9 full-slide images and current
speaker notes.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
SLIDES = ROOT / "docs" / "deck" / "slides"
OUT = ROOT / "docs" / "deck" / "Kampung SG-Project Introduction Deck-TheTwoGuys.pptx"

NOTES = [
    "Kampung SG is now a five-part campaign. Older residents are not problems "
    "to solve; their knowledge moves the story and visibly reshapes an "
    "enterable HDB estate. Every small act grows the kampung.",
    "The official challenge brief says that by 2030 nearly one in four "
    "Singapore citizens will be aged 65 and above. We respond with dignity, "
    "purpose, independence and social bonds: contributors rather than patients.",
    "The player moves from Y's flat through three sequential chapters and The "
    "Last Door, then keeps exploring. Homes, a corridor, workshop, shops and "
    "halls are enterable. The full story asks for three helpers and five "
    "invitees. There is no timer, energy bar or failure state.",
    "KampungMind is the headline AI feature. Each NPC has reviewed personality, "
    "role, expertise, knowledge, memory rules and authored intents. "
    "Deterministic scoring selects context-appropriate dialogue; stable IDs "
    "break ties and no text is generated at runtime.",
    "Our first CodeBuddy run failed its gate, and the correction attempt hit "
    "429 credits exhausted. We preserved that evidence. Later AI-assisted work "
    "added the pure campaign reducer, KampungMind and expanded checks. A real "
    "OpenAI image-generation pass produced one visual target; human review "
    "rejected direct runtime use and translated its strongest idea into four "
    "verified code-drawn activity vignettes. We treat AI output as a pull "
    "request, not a deliverable, and we do not relabel this as Miora.",
    "Strict TypeScript, 75 unit tests, 60 production-browser checks and zero "
    "known vulnerabilities support the current build. KampungMind and every "
    "save remain in the browser: no account, backend, analytics or personal "
    "data. Audio is synthesized at runtime.",
    "We separate software evidence from impact claims. The campaign, "
    "persistence, inputs and browser paths are automated; human playtesting, "
    "real-device accessibility and any social outcome are not claimed. "
    "Kampung SG makes no medical, cognitive, memory or dementia claim.",
    "We are TheTwoGuys. The public repository contains the current code, "
    "prompts, failed-run record, claim guardrails and reproducible gate. The "
    "game is playable in a browser with no install or account. Every small act "
    "grows the kampung.",
]


def main() -> None:
    images = sorted(SLIDES.glob("slide-*.png"))
    if len(images) != len(NOTES):
        raise SystemExit(f"Expected {len(NOTES)} slide images in {SLIDES}; found {len(images)}")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    for image, note in zip(images, NOTES, strict=True):
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )
        slide.notes_slide.notes_text_frame.text = note

    presentation.save(str(OUT))
    print(f"Wrote {OUT.name} - {len(images)} slides, {len(NOTES)} speaker notes")


if __name__ == "__main__":
    main()
